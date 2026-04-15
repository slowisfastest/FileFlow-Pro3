#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
FileFlow Pro - File Format Conversion Tool
Apple-style design, supports PDF/Word/Excel/PPT/Image/OFD format conversion
"""

import os
import sys
import traceback

# ===========================================================
# CRITICAL DEBUG: First thing - catch ALL startup failures
# ===========================================================
def _pause_on_error(msg="Press Enter to exit..."):
    """Pause console so user can read error output."""
    try:
        if getattr(sys, 'frozen', False):
            print(f"\n{'='*60}", flush=True)
            print(msg, flush=True)
            input()
    except Exception:
        pass

if getattr(sys, 'frozen', False):
    print(f"[DIAG] sys.executable = {sys.executable}", flush=True)
    print(f"[DIAG] sys._MEIPASS = {getattr(sys, '_MEIPASS', 'N/A')}", flush=True)
    print(f"[DIAG] Python {sys.version}", flush=True)

# ===========================================================
# OCR engine imports - wrapped for safety in frozen EXE
# ===========================================================
try:
    import paddle
    HAS_PADDLE = True
    print(f"[DIAG] paddle OK: {getattr(paddle, '__version__', '?')}", flush=True)
except Exception as _e:
    HAS_PADDLE = False
    print(f"[DIAG] paddle FAILED: {_e}", flush=True)

try:
    import paddleocr
    from paddleocr import PaddleOCR
    HAS_PADDLEOCR = True
    print(f"[DIAG] paddleocr OK", flush=True)
except Exception as _e:
    HAS_PADDLEOCR = False
    print(f"[DIAG] paddleocr FAILED: {_e}", flush=True)

try:
    import easyocr
    HAS_EASYOCR = True
    print(f"[DIAG] easyocr OK", flush=True)
except Exception as _e:
    HAS_EASYOCR = False
    print(f"[DIAG] easyocr FAILED: {_e}", flush=True)

try:
    import torch
    import torchvision
    HAS_TORCH = True
    print(f"[DIAG] torch OK: {getattr(torch, '__version__', '?')}", flush=True)
except Exception as _e:
    HAS_TORCH = False
    print(f"[DIAG] torch FAILED: {_e}", flush=True)

try:
    import cv2
    import numpy
    HAS_CV2 = True
    print(f"[DIAG] cv2 OK: {getattr(cv2, '__version__', '?')}", flush=True)
except Exception as _e:
    HAS_CV2 = False
    print(f"[DIAG] cv2 FAILED: {_e}", flush=True)

try:
    import PIL
    HAS_PIL = True
    print(f"[DIAG] PIL OK", flush=True)
except Exception as _e:
    HAS_PIL = False
    print(f"[DIAG] PIL FAILED: {_e}", flush=True)

try:
    import scipy
    HAS_SCIPY = True
    print(f"[DIAG] scipy OK: {getattr(scipy, '__version__', '?')}", flush=True)
except Exception as _e:
    HAS_SCIPY = False
    print(f"[DIAG] scipy FAILED: {_e}", flush=True)

try:
    import skimage
    HAS_SKIMAGE = True
    print(f"[DIAG] skimage OK", flush=True)
except Exception as _e:
    HAS_SKIMAGE = False
    print(f"[DIAG] skimage FAILED: {_e}", flush=True)

try:
    import shapely
    HAS_SHAPELY = True
    print(f"[DIAG] shapely OK", flush=True)
except Exception as _e:
    HAS_SHAPELY = False
    print(f"[DIAG] shapely FAILED: {_e}", flush=True)

# Only run diagnostics in frozen EXE mode
if getattr(sys, 'frozen', False):
    print(f"[DIAG] All imports done. HAS: P={HAS_PADDLE} PO={HAS_PADDLEOCR} E={HAS_EASYOCR} T={HAS_TORCH} C={HAS_CV2} PIL={HAS_PIL} SP={HAS_SCIPY} SK={HAS_SKIMAGE} SH={HAS_SHAPELY}", flush=True)
    print("[DIAG] Starting GUI...", flush=True)
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from pathlib import Path
import threading
import time
import shutil
import tempfile
import traceback
import platform

# ===========================================================
# Cross-platform file selector
# ===========================================================

def select_files_native(title="选择文件", multiple=True, file_types=None):
    """跨平台文件选择器"""
    system = platform.system()
    
    # macOS 使用原生选择器
    if system == "Darwin":
        try:
            from Cocoa import NSOpenPanel, NSModalResponseOK
            from Foundation import NSURL
            
            panel = NSOpenPanel.openPanel()
            panel.setTitle_(title)
            panel.setAllowsMultipleSelection_(multiple)
            panel.setCanChooseFiles_(True)
            panel.setCanChooseDirectories_(False)
            
            if file_types:
                allowed_types = []
                for ext in file_types:
                    ext_clean = ext.lstrip('.').lower()
                    allowed_types.append(ext_clean)
                panel.setAllowedFileTypes_(allowed_types)
            
            result = panel.runModal()
            
            if result == NSModalResponseOK:
                urls = panel.URLs()
                paths = [str(url.path()) for url in urls]
                return paths
            return []
        except:
            pass  # 失败则使用 tkinter
    
    # Windows 和 Linux 使用 tkinter filedialog
    filetypes = []
    if file_types:
        ext_list = [f"*.{ext.lstrip('.')}" for ext in file_types]
        filetypes.append(("支持的文件", " ".join(ext_list)))
    filetypes.append(("所有文件", "*.*"))
    
    if multiple:
        return list(filedialog.askopenfilenames(title=title, filetypes=filetypes))
    else:
        result = filedialog.askopenfilename(title=title, filetypes=filetypes)
        return [result] if result else []

def select_folder_native(title="选择文件夹", initial_dir=None):
    """跨平台文件夹选择器"""
    system = platform.system()
    
    # macOS 使用原生选择器
    if system == "Darwin":
        try:
            from Cocoa import NSOpenPanel, NSModalResponseOK
            from Foundation import NSURL
            
            panel = NSOpenPanel.openPanel()
            panel.setTitle_(title)
            panel.setAllowsMultipleSelection_(False)
            panel.setCanChooseFiles_(False)
            panel.setCanChooseDirectories_(True)
            
            if initial_dir and os.path.exists(initial_dir):
                url = NSURL.fileURLWithPath_(initial_dir)
                panel.setDirectoryURL_(url)
            
            result = panel.runModal()
            
            if result == NSModalResponseOK:
                urls = panel.URLs()
                if urls:
                    return str(urls[0].path())
            return None
        except:
            pass  # 失败则使用 tkinter
    
    # Windows 和 Linux 使用 tkinter filedialog
    return filedialog.askdirectory(title=title, initialdir=initial_dir)

# ===========================================================
# Configuration and constants
# ===========================================================

C = {
    "bg":          "#F5F5F7",
    "surface":     "#FFFFFF", 
    "blue":        "#0051D5",
    "blue_light":  "#E5F1FF",
    "green":       "#248A3D",
    "red":         "#D70015",
    "text":        "#1C1C1E",
    "text2":       "#3C3C43",
    "placeholder": "#8E8E93",
    "border":      "#D1D1D6",
}

# 支持的文件扩展名
SUPPORTED_EXTS = ['pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ofd', 
                  'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'tif', 'webp', 
                  'txt', 'csv', 'html', 'htm', 'md']

CONVERT_OPTIONS = {
    "pdf":  ["Word (.docx)", "Word-OCR (.docx)", "Excel (.xlsx)", "PPT (.pptx)", "图片 (.png)", "文本 (.txt)"],
    "docx": ["PDF (.pdf)", "文本 (.txt)", "HTML (.html)"],
    "xlsx": ["PDF (.pdf)", "CSV (.csv)", "文本 (.txt)"],
    "pptx": ["PDF (.pdf)", "图片 (.png)"],
    "ofd":  ["PDF (.pdf)", "Word (.docx)", "文本 (.txt)"],
    "png":  ["PDF (.pdf)", "JPG (.jpg)", "WebP (.webp)"],
    "jpg":  ["PDF (.pdf)", "PNG (.png)", "WebP (.webp)"],
    "jpeg": ["PDF (.pdf)", "PNG (.png)", "WebP (.webp)"],
    "txt":  ["PDF (.pdf)", "Word (.docx)"],
}

# 格式到扩展名的映射
EXT_MAP = {
    "Word (.docx)": ".docx",
    "Word-OCR (.docx)": ".docx",
    "Excel (.xlsx)": ".xlsx",
    "PPT (.pptx)": ".pptx",
    "图片 (.png)": ".png",
    "JPG (.jpg)": ".jpg",
    "PNG (.png)": ".png",
    "PDF (.pdf)": ".pdf",
    "文本 (.txt)": ".txt",
    "HTML (.html)": ".html",
    "CSV (.csv)": ".csv",
    "WebP (.webp)": ".webp",
}

# ===========================================================
# Conversion engine
# ===========================================================

class ConversionEngine:
    """核心转换逻辑"""
    
    @staticmethod
    def pdf_to_word(src: str, dst: str, log_func=None, use_ocr=False):
        """PDF 转 Word - 支持原生文本和 OCR 扫描件"""
        try:
            import fitz  # PyMuPDF
            from docx import Document
            from docx.shared import Pt, Inches
            from PIL import Image
            import io
            
            if log_func:
                log_func("  分析 PDF 内容...")
            
            # 打开 PDF
            doc = fitz.open(src)
            word_doc = Document()
            
            # 检测是否主要是扫描件/图片
            is_scan = use_ocr or ConversionEngine._is_scanned_pdf(doc, log_func)
            
            if is_scan:
                if log_func:
                    log_func("  检测到扫描件/图片 PDF，启用 OCR 识别...")
                return ConversionEngine._pdf_to_word_with_ocr(doc, word_doc, dst, log_func)
            else:
                if log_func:
                    log_func("  检测到文本 PDF，使用原生转换...")
                return ConversionEngine._pdf_to_word_native(src, dst, log_func)
                
        except Exception as e:
            raise RuntimeError(f"PDF转Word失败: {e}")
    
    @staticmethod
    def _is_scanned_pdf(doc, log_func=None):
        """检测 PDF 是否主要是扫描件/图片"""
        try:
            total_pages = len(doc)
            if total_pages == 0:
                return False
            
            text_pages = 0
            image_pages = 0
            empty_pages = 0
            
            # 采样检查前 3 页
            sample_pages = min(3, total_pages)
            for i in range(sample_pages):
                page = doc[i]
                
                # 检查文本量
                text = page.get_text().strip()
                text_length = len(text)
                
                # 检查图片数量（嵌入的图片对象）
                image_list = page.get_images()
                
                # 检查页面是否有绘制命令（用于判断是否是纯图像页面）
                # 扫描件通常有大量绘制命令但无文本
                xrefs = len(image_list)
                
                # 改进的检测逻辑：
                # 1. 如果文本很少（<50字符）但有嵌入图片，可能是扫描件
                # 2. 如果文本很少且页面看起来是图片类型，可能是扫描件
                # 3. 如果文本为0或极少，很可能是扫描件
                
                if text_length == 0:
                    # 完全无文本，很可能是扫描件
                    empty_pages += 1
                    image_pages += 1
                elif text_length < 50 and xrefs > 0:
                    # 极少文本但有图片对象
                    image_pages += 1
                elif text_length < 100 and xrefs > 0:
                    # 较少文本且有图片对象
                    image_pages += 1
                elif text_length > 200:
                    # 有足够文本，认为是文本页面
                    text_pages += 1
                else:
                    # 边界情况：检查页面内容类型
                    # 如果页面主要是图像数据，认为是扫描件
                    if xrefs > 0 or text_length < 30:
                        image_pages += 1
                    else:
                        text_pages += 1
            
            if log_func:
                log_func(f"  页面分析: 文本页={text_pages}, 图片/扫描页={image_pages}, 空白页={empty_pages}")
            
            # 如果采样页大部分是图片型或空白，判定为扫描件
            return image_pages > text_pages or empty_pages >= sample_pages // 2
            
        except Exception as e:
            if log_func:
                log_func(f"  扫描检测出错，默认使用原生转换: {e}")
            return False
    
    @staticmethod
    def _pdf_to_word_native(src: str, dst: str, log_func=None):
        """使用 pdf2docx 进行原生转换"""
        try:
            from pdf2docx import Converter
            if log_func:
                log_func("  使用 pdf2docx 转换...")
            cv = Converter(src)
            cv.convert(dst, start=0, end=None)
            cv.close()
            return True
        except Exception as e:
            raise RuntimeError(f"原生转换失败: {e}")
    
    @staticmethod
    def _get_ocr_model_path():
        """获取 OCR 模型路径（支持打包后的 EXE 环境）"""
        import sys
        
        # 检测是否在 PyInstaller 打包环境中
        if getattr(sys, 'frozen', False):
            # 打包后的 EXE 环境
            base_path = Path(sys._MEIPASS)
        else:
            # 开发环境
            base_path = Path(__file__).parent
        
        models_dir = base_path / "ocr_models"
        return models_dir if models_dir.exists() else None
    
    @staticmethod
    def _pdf_to_word_with_ocr(doc, word_doc, dst: str, log_func=None):
        """Convert scanned PDF to editable Word using OCR"""
        try:
            import fitz  # ensure fitz is available in this scope
            from PIL import Image
            import io
            import numpy as np
            
            # 尝试导入 OCR 库 - 优先使用 PaddleOCR (SOTA 模型)
            ocr_available = False
            ocr_reader = None
            ocr_engine = None  # 初始化引擎类型变量
            
            # 获取离线模型路径
            model_path = ConversionEngine._get_ocr_model_path()
            
            # 优先使用 PaddleOCR-VL (2025-2026 SOTA, OmniDocBench 92.86分)
            try:
                from paddleocr import PaddleOCR
                if log_func:
                    log_func("  加载 PaddleOCR 引擎...")
                
                # 配置 PaddleOCR 参数
                paddle_kwargs = {
                    'use_angle_cls': True,  # 方向分类
                    'lang': 'ch',  # 中文模型
                    'use_gpu': False,  # CPU 运行
                }
                
                # 尝试添加 show_log 参数（新版本支持）
                try:
                    import inspect
                    sig = inspect.signature(PaddleOCR.__init__)
                    if 'show_log' in sig.parameters:
                        paddle_kwargs['show_log'] = False  # 关闭冗余日志
                except:
                    pass
                
                # 如果存在离线模型，使用离线模型路径
                if model_path:
                    paddle_dir = model_path / "paddleocr"
                    if paddle_dir.exists():
                        if log_func:
                            log_func(f"  使用离线模型: {paddle_dir}")
                        # 指定模型路径
                        det_model_dir = paddle_dir / "ch_PP-OCRv4_det_infer"
                        rec_model_dir = paddle_dir / "ch_PP-OCRv4_rec_infer"
                        cls_model_dir = paddle_dir / "ch_ppocr_mobile_v2.0_cls_infer"
                        
                        if det_model_dir.exists():
                            paddle_kwargs['det_model_dir'] = str(det_model_dir)
                        if rec_model_dir.exists():
                            paddle_kwargs['rec_model_dir'] = str(rec_model_dir)
                        if cls_model_dir.exists():
                            paddle_kwargs['cls_model_dir'] = str(cls_model_dir)
                
                # 初始化 PaddleOCR
                ocr_reader = PaddleOCR(**paddle_kwargs)
                ocr_available = True
                ocr_engine = "paddleocr"
                if log_func:
                    log_func("  [OK] PaddleOCR loaded")
            except Exception as e:
                if log_func:
                    log_func(f"  PaddleOCR 加载失败: {e}")
                    log_func("  尝试 EasyOCR...")
            
            # 备选使用 EasyOCR
            if not ocr_available:
                try:
                    import easyocr
                    import os
                    if log_func:
                        log_func("  加载 EasyOCR 引擎...")
                    
                    # 检查是否存在离线模型
                    easyocr_model_dir = None
                    if model_path:
                        easyocr_dir = model_path / "easyocr"
                        if easyocr_dir.exists():
                            easyocr_model_dir = str(easyocr_dir)
                            if log_func:
                                log_func(f"  使用离线 EasyOCR 模型: {easyocr_model_dir}")
                            # 设置环境变量指定模型路径
                            os.environ['EASYOCR_MODULE_PATH'] = easyocr_model_dir
                    
                    ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, model_storage_directory=easyocr_model_dir)
                    ocr_available = True
                    ocr_engine = "easyocr"
                    if log_func:
                        log_func("  [OK] EasyOCR loaded")
                except Exception as e:
                    if log_func:
                        log_func(f"  EasyOCR 加载失败: {e}")
                        log_func("  尝试 Tesseract...")
            
            # 最后备选使用 pytesseract
            if not ocr_available:
                try:
                    import pytesseract
                    from PIL import Image
                    if log_func:
                        log_func("  使用 Tesseract OCR 引擎...")
                    ocr_available = True
                    ocr_engine = "tesseract"
                    if log_func:
                        log_func("  [OK] Tesseract loaded")
                except Exception as e:
                    if log_func:
                        log_func(f"  Tesseract 加载失败: {e}")
            
            if not ocr_available:
                raise RuntimeError(
                    "No OCR engine available. Please ensure PaddleOCR or EasyOCR is properly installed.\n"
                    "Error details will be logged above."
                )
            
            total_pages = len(doc)
            
            for page_num in range(total_pages):
                if log_func:
                    log_func(f"  处理第 {page_num + 1}/{total_pages} 页...")
                
                page = doc[page_num]
                
                # 渲染页面为图片
                mat = fitz.Matrix(2, 2)  # 2x 缩放以获得更好识别效果
                pix = page.get_pixmap(matrix=mat)
                
                # 转换为 PIL Image
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                # 执行 OCR
                if ocr_engine == "paddleocr":
                    # PaddleOCR-VL 识别 (SOTA 模型)
                    import cv2
                    
                    # 将 PIL Image 转换为 numpy array
                    img_array = np.array(img)
                    
                    # PaddleOCR 需要 BGR 格式
                    if len(img_array.shape) == 3:
                        img_array = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
                    
                    # 执行 OCR
                    result = ocr_reader.ocr(img_array, cls=True)
                    
                    # 按行组织文本
                    lines = []
                    current_y = None
                    current_line = []
                    
                    if result and result[0]:
                        for line_data in result[0]:
                            if line_data:
                                bbox = line_data[0]  # 边界框
                                text = line_data[1][0]  # 文本内容
                                conf = line_data[1][1]  # 置信度
                                
                                if conf < 0.3:  # 过滤低置信度
                                    continue
                                
                                # 计算文本块的中心 Y 坐标
                                y_center = (bbox[0][1] + bbox[2][1]) / 2
                                
                                if current_y is None or abs(y_center - current_y) < 20:
                                    current_line.append(text)
                                    current_y = y_center
                                else:
                                    if current_line:
                                        lines.append(" ".join(current_line))
                                    current_line = [text]
                                    current_y = y_center
                        
                        if current_line:
                            lines.append(" ".join(current_line))
                    
                    # 添加到 Word 文档
                    for line in lines:
                        if line.strip():
                            word_doc.add_paragraph(line)
                    
                elif ocr_engine == "easyocr":
                    # EasyOCR 识别
                    result = ocr_reader.readtext(np.array(img))
                    
                    # 按行组织文本
                    lines = []
                    current_y = None
                    current_line = []
                    
                    for (bbox, text, conf) in result:
                        if conf < 0.3:  # 过滤低置信度
                            continue
                        
                        # 计算文本块的中心 Y 坐标
                        y_center = (bbox[0][1] + bbox[2][1]) / 2
                        
                        if current_y is None or abs(y_center - current_y) < 20:
                            current_line.append(text)
                            current_y = y_center
                        else:
                            if current_line:
                                lines.append(" ".join(current_line))
                            current_line = [text]
                            current_y = y_center
                    
                    if current_line:
                        lines.append(" ".join(current_line))
                    
                    # 添加到 Word 文档
                    for line in lines:
                        if line.strip():
                            word_doc.add_paragraph(line)
                    
                else:
                    # Tesseract OCR
                    text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    
                    # 添加到 Word 文档
                    for line in text.split('\n'):
                        if line.strip():
                            word_doc.add_paragraph(line)
                
                # 页之间添加空行
                if page_num < total_pages - 1:
                    word_doc.add_paragraph()
            
            # 保存 Word 文档
            word_doc.save(dst)
            
            if log_func:
                log_func(f"  [OK] OCR done, {total_pages} pages total")
            
            return True
            
        except Exception as e:
            raise RuntimeError(f"OCR转换失败: {e}")
    
    @staticmethod
    def pdf_to_excel(src: str, dst: str, log_func=None):
        """PDF 转 Excel - 使用多种策略提高表格识别准确率"""
        try:
            import pdfplumber
            import openpyxl
            from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
            import re
            
            if log_func:
                log_func("  提取PDF表格数据...")
            
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            
            # 样式定义
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font_white = Font(bold=True, size=11, color="FFFFFF")
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            with pdfplumber.open(src) as pdf:
                row_cursor = 1
                for page_num, page in enumerate(pdf.pages, 1):
                    if log_func:
                        log_func(f"  处理第 {page_num} 页...")
                    
                    # 策略1: 尝试线条模式提取表格
                    tables = page.extract_tables({
                        "vertical_strategy": "lines",
                        "horizontal_strategy": "lines",
                        "intersection_y_tolerance": 5,
                        "intersection_x_tolerance": 5,
                    })
                    
                    # 策略2: 如果没找到，尝试文本对齐模式
                    if not tables or all(len(t) == 0 for t in tables):
                        tables = page.extract_tables({
                            "vertical_strategy": "text",
                            "horizontal_strategy": "text",
                            "snap_tolerance": 3,
                            "join_tolerance": 3,
                            "edge_min_length": 10,
                        })
                    
                    # 策略3: 显式模式
                    if not tables or all(len(t) == 0 for t in tables):
                        tables = page.extract_tables({
                            "vertical_strategy": "explicit",
                            "horizontal_strategy": "explicit",
                        })
                    
                    if tables:
                        for table_idx, table in enumerate(tables):
                            if not table:
                                continue
                            
                            if log_func:
                                log_func(f"  发现表格 {table_idx + 1}，{len(table)} 行...")
                            
                            # 确定表格的最大列数
                            max_cols = max(len(row) for row in table) if table else 0
                            
                            for row_idx, row in enumerate(table):
                                if not row:
                                    continue
                                
                                # 清理和标准化单元格数据
                                cleaned_row = []
                                for cell in row:
                                    if cell is None:
                                        cleaned_row.append("")
                                    else:
                                        # 清理特殊字符，保留基本格式
                                        cell_text = str(cell).strip()
                                        # 保留换行符但清理其他控制字符
                                        cell_text = cell_text.replace('\r', '').replace('\t', ' ')
                                        cleaned_row.append(cell_text)
                                
                                # 补齐缺失的列
                                while len(cleaned_row) < max_cols:
                                    cleaned_row.append("")
                                
                                # 写入行
                                for col_idx, value in enumerate(cleaned_row):
                                    cell = ws.cell(row=row_cursor, column=col_idx + 1, value=value)
                                    cell.border = thin_border
                                    cell.alignment = Alignment(vertical='center', wrap_text=True)
                                    
                                    # 第一行设为表头样式（假设第一行是表头）
                                    if row_idx == 0:
                                        cell.font = header_font_white
                                        cell.fill = header_fill
                                        cell.alignment = Alignment(horizontal='center', vertical='center')
                                
                                row_cursor += 1
                            
                            # 表格之间空一行
                            row_cursor += 1
                    else:
                        # 策略4: 智能文本解析 - 尝试识别表格结构
                        text = page.extract_text() or ""
                        if text.strip():
                            lines = text.split('\n')
                            table_data = []
                            current_row = []
                            
                            for line in lines:
                                line = line.strip()
                                if not line:
                                    continue
                                
                                # 检测是否是数据行（包含数字、日期、货币等模式）
                                # 或者包含多个空格分隔的字段
                                is_data_row = False
                                
                                # 检查是否包含表格特征（多个空格分隔、制表符、竖线等）
                                if '|' in line or '\t' in line:
                                    is_data_row = True
                                    parts = [p.strip() for p in re.split(r'[|\t]', line) if p.strip()]
                                elif re.search(r'\d{4}[-/]\d{2}[-/]\d{2}', line):  # 日期
                                    is_data_row = True
                                    parts = [p.strip() for p in line.split() if p.strip()]
                                elif re.search(r'\d+\.\d{2}', line):  # 金额格式
                                    is_data_row = True
                                    parts = [p.strip() for p in line.split() if p.strip()]
                                elif len(line.split()) >= 4:  # 多个字段
                                    is_data_row = True
                                    parts = [p.strip() for p in line.split() if p.strip()]
                                else:
                                    parts = [line]
                                
                                if is_data_row and len(parts) >= 2:
                                    for col_idx, value in enumerate(parts):
                                        cell = ws.cell(row=row_cursor, column=col_idx + 1, value=value)
                                        cell.border = thin_border
                                        cell.alignment = Alignment(vertical='center', wrap_text=True)
                                    row_cursor += 1
                                else:
                                    # 普通文本行，作为单列
                                    ws.cell(row=row_cursor, column=1, value=line)
                                    row_cursor += 1
                            
                            row_cursor += 1
            
            # 自动调整列宽
            if log_func:
                log_func("  调整格式...")
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 60)
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # 冻结首行
            if ws.max_row > 1:
                ws.freeze_panes = 'A2'
            
            wb.save(dst)
            return True
            
        except Exception as e:
            raise RuntimeError(f"PDF转Excel失败: {e}")
    
    @staticmethod
    def pdf_to_pptx(src: str, dst: str, log_func=None):
        """PDF 转 PPT"""
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Inches
            
            if log_func:
                log_func("  提取PDF页面...")
            
            doc = fitz.open(src)
            prs = Presentation()
            blank = prs.slide_layouts[6]
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            
            for i, page in enumerate(doc):
                if log_func:
                    log_func(f"  处理第 {i+1}/{len(doc)} 页...")
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                slide = prs.slides.add_slide(blank)
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                    f.write(img_bytes)
                    tmp = f.name
                slide.shapes.add_picture(tmp, 0, 0, slide_w, slide_h)
                os.unlink(tmp)
            
            prs.save(dst)
            return True
        except Exception as e:
            raise RuntimeError(f"PDF转PPT失败: {e}")
    
    @staticmethod
    def pdf_to_image(src: str, dst: str, fmt="png", log_func=None):
        """PDF 转图片"""
        try:
            import fitz
            from PIL import Image
            
            if log_func:
                log_func("  渲染PDF...")
            
            doc = fitz.open(src)
            if len(doc) == 1:
                page = doc[0]
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                pix.save(dst)
            else:
                # 多页合并
                base, _ = os.path.splitext(dst)
                paths = []
                for i, page in enumerate(doc):
                    if log_func:
                        log_func(f"  处理第 {i+1}/{len(doc)} 页...")
                    mat = fitz.Matrix(2, 2)
                    pix = page.get_pixmap(matrix=mat)
                    p = f"{base}_page{i+1}.{fmt}"
                    pix.save(p)
                    paths.append(p)
                
                # 垂直拼接
                imgs = [Image.open(p) for p in paths]
                total_h = sum(im.height for im in imgs)
                max_w = max(im.width for im in imgs)
                merged = Image.new("RGB", (max_w, total_h), (255, 255, 255))
                y = 0
                for im in imgs:
                    x_offset = (max_w - im.width) // 2
                    merged.paste(im, (x_offset, y))
                    y += im.height
                merged.save(dst)
                for p in paths:
                    os.unlink(p)
            return True
        except Exception as e:
            raise RuntimeError(f"PDF转图片失败: {e}")
    
    @staticmethod
    def pdf_to_txt(src: str, dst: str, log_func=None):
        """PDF 转文本"""
        try:
            import pdfplumber
            
            if log_func:
                log_func("  提取文本...")
            
            lines = []
            with pdfplumber.open(src) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    if log_func:
                        log_func(f"  处理第 {i}/{len(pdf.pages)} 页...")
                    text = page.extract_text() or ""
                    lines.append(f"=== 第 {i} 页 ===")
                    lines.append(text)
                    lines.append("")
            
            with open(dst, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))
            return True
        except Exception as e:
            raise RuntimeError(f"PDF转文本失败: {e}")
    
    @staticmethod
    def word_to_pdf(src: str, dst: str, log_func=None):
        """Word 转 PDF - 使用 python-docx + reportlab"""
        try:
            if log_func:
                log_func("  读取 Word 文档...")
            
            from docx import Document
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.lib import colors
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            # 注册中文字体
            font_name = "Helvetica"
            font_paths = [
                # macOS
                "/System/Library/Fonts/PingFang.ttc",
                "/System/Library/Fonts/STHeiti Light.ttc",
                # Windows
                "C:/Windows/Fonts/msyh.ttc",
                "C:/Windows/Fonts/simsun.ttc",
                "C:/Windows/Fonts/simhei.ttf",
                "C:/Windows/Fonts/simkai.ttf",
            ]
            for fp in font_paths:
                if os.path.exists(fp):
                    try:
                        pdfmetrics.registerFont(TTFont("CJK", fp))
                        font_name = "CJK"
                        break
                    except:
                        pass
            
            if log_func:
                log_func("  创建 PDF...")
            
            # 读取 Word 文档
            doc = Document(src)
            
            # 创建 PDF
            pdf_doc = SimpleDocTemplate(dst, pagesize=A4,
                                        leftMargin=20*mm, rightMargin=20*mm,
                                        topMargin=20*mm, bottomMargin=20*mm)
            
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                "Title",
                parent=styles["Heading1"],
                fontName=font_name,
                fontSize=18,
                leading=24,
                spaceAfter=12,
            )
            body_style = ParagraphStyle(
                "Body",
                parent=styles["Normal"],
                fontName=font_name,
                fontSize=11,
                leading=16,
                spaceAfter=6,
            )
            
            story = []
            
            # 正确迭代所有段落和表格（保持顺序）
            from docx.oxml.ns import qn
            para_idx = 0
            tbl_idx = 0
            para_objs = doc.paragraphs
            tbl_objs = doc.tables
            
            for element in doc.element.body:
                tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
                if tag == 'p':  # 段落
                    if para_idx < len(para_objs):
                        para = para_objs[para_idx]
                        para_idx += 1
                        if para.text.strip():
                            text = para.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                            # 根据段落样式选择对应样式
                            if para.style and para.style.name and 'Heading' in para.style.name:
                                story.append(Paragraph(text, title_style))
                            else:
                                story.append(Paragraph(text, body_style))
                elif tag == 'tbl':  # 表格
                    if tbl_idx < len(tbl_objs):
                        table = tbl_objs[tbl_idx]
                        tbl_idx += 1
                        data = []
                        for row in table.rows:
                            row_data = [cell.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;") for cell in row.cells]
                            data.append(row_data)
                        
                        if data:
                            t = Table(data)
                            t.setStyle(TableStyle([
                                ('FONTNAME', (0, 0), (-1, -1), font_name),
                                ('FONTSIZE', (0, 0), (-1, -1), 10),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#E8E8E8")),
                            ]))
                            story.append(t)
                            story.append(Spacer(1, 6))
            
            # 如果遍历后没有内容（极少数情况），退回直接读段落
            if not story:
                for para in doc.paragraphs:
                    if para.text.strip():
                        text = para.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                        story.append(Paragraph(text, body_style))
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            raise RuntimeError(f"Word转PDF失败: {e}")
    
    @staticmethod
    def word_to_txt(src: str, dst: str, log_func=None):
        """Word 转文本"""
        try:
            from docx import Document
            
            if log_func:
                log_func("  提取文本...")
            
            doc = Document(src)
            text = "\n".join(p.text for p in doc.paragraphs)
            with open(dst, "w", encoding="utf-8") as f:
                f.write(text)
            return True
        except Exception as e:
            raise RuntimeError(f"Word转文本失败: {e}")
    
    @staticmethod
    def word_to_html(src: str, dst: str, log_func=None):
        """Word 转 HTML"""
        try:
            from docx import Document
            from docx.oxml.ns import qn
            import html as html_mod
            
            if log_func:
                log_func("  转换为 HTML...")
            
            def escape_html(text):
                """转义 HTML 特殊字符，防止注入"""
                return html_mod.escape(str(text))
            
            doc = Document(src)
            html_lines = [
                "<!DOCTYPE html>",
                "<html lang='zh-CN'>",
                "<head>",
                "<meta charset='UTF-8'>",
                "<meta name='viewport' content='width=device-width, initial-scale=1.0'>",
                "<style>",
                "body { font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; line-height: 1.6; }",
                "h1 { color: #333; } h2 { color: #444; } h3 { color: #555; }",
                "p { margin: 8px 0; }",
                "table { border-collapse: collapse; width: 100%; margin: 10px 0; }",
                "td, th { border: 1px solid #ddd; padding: 8px; text-align: left; }",
                "th { background-color: #f2f2f2; }",
                "</style>",
                "</head>",
                "<body>",
            ]
            
            # 按文档顺序处理段落和表格
            para_idx = 0
            tbl_idx = 0
            para_objs = doc.paragraphs
            tbl_objs = doc.tables
            
            for element in doc.element.body:
                tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
                
                if tag == 'p':
                    if para_idx < len(para_objs):
                        para = para_objs[para_idx]
                        para_idx += 1
                        text = para.text.strip()
                        if not text:
                            html_lines.append("<p><br></p>")
                            continue
                        
                        # 根据段落样式选择 HTML 标签
                        if para.style and para.style.name:
                            style_name = para.style.name
                            if 'Heading 1' in style_name or 'Title' in style_name:
                                html_lines.append(f"<h1>{escape_html(text)}</h1>")
                            elif 'Heading 2' in style_name:
                                html_lines.append(f"<h2>{escape_html(text)}</h2>")
                            elif 'Heading 3' in style_name:
                                html_lines.append(f"<h3>{escape_html(text)}</h3>")
                            else:
                                html_lines.append(f"<p>{escape_html(text)}</p>")
                        else:
                            html_lines.append(f"<p>{escape_html(text)}</p>")
                
                elif tag == 'tbl':
                    if tbl_idx < len(tbl_objs):
                        table = tbl_objs[tbl_idx]
                        tbl_idx += 1
                        html_lines.append("<table>")
                        for row_idx, row in enumerate(table.rows):
                            tag_name = "th" if row_idx == 0 else "td"
                            html_lines.append("<tr>")
                            for cell in row.cells:
                                html_lines.append(f"<{tag_name}>{escape_html(cell.text)}</{tag_name}>")
                            html_lines.append("</tr>")
                        html_lines.append("</table>")
            
            html_lines.extend(["</body>", "</html>"])
            
            with open(dst, "w", encoding="utf-8") as f:
                f.write("\n".join(html_lines))
            
            return True
        except Exception as e:
            raise RuntimeError(f"Word转HTML失败: {e}")
    
    @staticmethod
    def excel_to_pdf(src: str, dst: str, log_func=None):
        """Excel 转 PDF - 使用 openpyxl + reportlab"""
        try:
            import openpyxl
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
            from reportlab.lib import colors
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.lib.units import mm
            
            if log_func:
                log_func("  读取 Excel...")
            
            # 注册中文字体
            font_name = "Helvetica"
            font_paths = [
                # macOS
                "/System/Library/Fonts/PingFang.ttc",
                "/System/Library/Fonts/STHeiti Light.ttc",
                # Windows
                "C:/Windows/Fonts/msyh.ttc",
                "C:/Windows/Fonts/simsun.ttc",
                "C:/Windows/Fonts/simhei.ttf",
                "C:/Windows/Fonts/simkai.ttf",
            ]
            for fp in font_paths:
                if os.path.exists(fp):
                    try:
                        pdfmetrics.registerFont(TTFont("CJK", fp))
                        font_name = "CJK"
                        break
                    except:
                        pass
            
            # 读取 Excel
            wb = openpyxl.load_workbook(src, data_only=True)
            ws = wb.active
            
            if log_func:
                log_func("  创建 PDF...")
            
            # 创建 PDF（横向以容纳更多列）
            pdf_doc = SimpleDocTemplate(dst, pagesize=landscape(A4),
                                        leftMargin=10*mm, rightMargin=10*mm,
                                        topMargin=10*mm, bottomMargin=10*mm)
            
            # 提取数据
            data = []
            max_cols = 0
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                data.append(row_data)
                max_cols = max(max_cols, len(row_data))
            
            # 确保每行数据长度一致
            for row in data:
                while len(row) < max_cols:
                    row.append("")
            
            # 创建表格
            if data:
                # 计算列宽
                col_widths = []
                for col_idx in range(max_cols):
                    max_len = 0
                    for row in data:
                        if col_idx < len(row):
                            max_len = max(max_len, len(str(row[col_idx])))
                    # 根据内容长度设置列宽（字符数 * 6mm）
                    col_widths.append(min(max(max_len * 3, 20), 80) * mm)
                
                t = Table(data, colWidths=col_widths, repeatRows=1)
                t.setStyle(TableStyle([
                    ('FONTNAME', (0, 0), (-1, -1), font_name),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor("#F2F2F2")]),
                    ('LEFTPADDING', (0, 0), (-1, -1), 4),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 4),
                    ('TOPPADDING', (0, 0), (-1, -1), 4),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                ]))
                
                pdf_doc.build([t])
            else:
                # 空表格
                pdf_doc.build([Paragraph("空表格", ParagraphStyle("Empty", fontName=font_name, fontSize=12))])
            
            return True
        except Exception as e:
            raise RuntimeError(f"Excel转PDF失败: {e}")
    
    @staticmethod
    def excel_to_csv(src: str, dst: str, log_func=None):
        """Excel 转 CSV"""
        try:
            import openpyxl
            import csv
            
            if log_func:
                log_func("  转换中...")
            
            wb = openpyxl.load_workbook(src, data_only=True)
            ws = wb.active
            with open(dst, "w", newline="", encoding="utf-8-sig") as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow([v if v is not None else "" for v in row])
            return True
        except Exception as e:
            raise RuntimeError(f"Excel转CSV失败: {e}")
    
    @staticmethod
    def excel_to_txt(src: str, dst: str, log_func=None):
        """Excel 转文本（制表符分隔，可读性好）"""
        try:
            import openpyxl
            
            if log_func:
                log_func("  转换中...")
            
            wb = openpyxl.load_workbook(src, data_only=True)
            lines = []
            for ws_name in wb.sheetnames:
                ws = wb[ws_name]
                if len(wb.sheetnames) > 1:
                    lines.append(f"=== 工作表: {ws_name} ===")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    lines.append("\t".join(cells))
                lines.append("")
            
            with open(dst, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))
            return True
        except Exception as e:
            raise RuntimeError(f"Excel转文本失败: {e}")
    
    @staticmethod
    def pptx_to_pdf(src: str, dst: str, log_func=None):
        """PPT 转 PDF"""
        try:
            soffice_paths = [
                # macOS
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                # Windows
                "C:/Program Files/LibreOffice/program/soffice.exe",
                "C:/Program Files (x86)/LibreOffice/program/soffice.exe",
                # Linux/通用
                "soffice",
                "libreoffice"
            ]
            soffice = None
            for path in soffice_paths:
                if shutil.which(path) or os.path.exists(path):
                    soffice = path
                    break
            
            if soffice:
                import subprocess
                if log_func:
                    log_func("  使用 LibreOffice 转换...")
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf", "--outdir",
                     os.path.dirname(dst), src],
                    capture_output=True, timeout=120
                )
                lo_out = os.path.join(os.path.dirname(dst), Path(src).stem + ".pdf")
                if os.path.exists(lo_out) and os.path.abspath(lo_out) != os.path.abspath(dst):
                    shutil.move(lo_out, dst)
                if os.path.exists(dst):
                    return True
            
            raise RuntimeError("需要安装 LibreOffice")
        except Exception as e:
            raise RuntimeError(f"PPT转PDF失败: {e}")
    
    @staticmethod
    def pptx_to_image(src: str, dst: str, log_func=None):
        """PPT 转图片"""
        try:
            import fitz
            
            if log_func:
                log_func("  先转PDF再转图片...")
            
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                tmp_pdf = f.name
            
            ConversionEngine.pptx_to_pdf(src, tmp_pdf, log_func)
            ConversionEngine.pdf_to_image(tmp_pdf, dst, fmt="png", log_func=log_func)
            os.unlink(tmp_pdf)
            return True
        except Exception as e:
            raise RuntimeError(f"PPT转图片失败: {e}")
    
    @staticmethod
    def image_to_pdf(src: str, dst: str, log_func=None):
        """图片 转 PDF"""
        try:
            from PIL import Image
            
            if log_func:
                log_func("  转换中...")
            
            img = Image.open(src)
            if img.mode in ("RGBA", "P", "LA"):
                img = img.convert("RGB")
            img.save(dst, "PDF", resolution=100)
            return True
        except Exception as e:
            raise RuntimeError(f"图片转PDF失败: {e}")
    
    @staticmethod
    def image_to_image(src: str, dst: str, log_func=None):
        """图片 转 图片（格式转换）"""
        try:
            from PIL import Image
            
            if log_func:
                log_func("  转换格式...")
            
            img = Image.open(src)
            ext = Path(dst).suffix.lower()
            if ext in (".jpg", ".jpeg") and img.mode in ("RGBA", "P", "LA"):
                img = img.convert("RGB")
            img.save(dst)
            return True
        except Exception as e:
            raise RuntimeError(f"图片格式转换失败: {e}")
    
    @staticmethod
    def txt_to_pdf(src: str, dst: str, log_func=None):
        """文本 转 PDF"""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            if log_func:
                log_func("  创建PDF...")
            
            # 尝试注册中文字体
            font_name = "Helvetica"
            font_paths = [
                # macOS
                "/System/Library/Fonts/PingFang.ttc",
                "/System/Library/Fonts/STHeiti Light.ttc",
                # Windows
                "C:/Windows/Fonts/msyh.ttc",
                "C:/Windows/Fonts/simsun.ttc",
                "C:/Windows/Fonts/simhei.ttf",
                "C:/Windows/Fonts/simkai.ttf",
            ]
            for fp in font_paths:
                if os.path.exists(fp):
                    try:
                        pdfmetrics.registerFont(TTFont("CJK", fp))
                        font_name = "CJK"
                        break
                    except:
                        pass
            
            with open(src, "r", encoding="utf-8") as f:
                content = f.read()
            
            styles = getSampleStyleSheet()
            body_style = ParagraphStyle(
                "Body",
                parent=styles["Normal"],
                fontName=font_name,
                fontSize=11,
                leading=18,
                spaceAfter=6,
            )
            
            doc = SimpleDocTemplate(dst, pagesize=A4,
                                    leftMargin=20*mm, rightMargin=20*mm,
                                    topMargin=20*mm, bottomMargin=20*mm)
            story = []
            for line in content.split("\n"):
                safe_line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;") or " "
                story.append(Paragraph(safe_line, body_style))
            
            doc.build(story)
            return True
        except Exception as e:
            raise RuntimeError(f"文本转PDF失败: {e}")
    
    @staticmethod
    def txt_to_word(src: str, dst: str, log_func=None):
        """文本 转 Word"""
        try:
            from docx import Document
            
            if log_func:
                log_func("  创建Word文档...")
            
            doc = Document()
            with open(src, "r", encoding="utf-8") as f:
                for line in f:
                    doc.add_paragraph(line.rstrip())
            doc.save(dst)
            return True
        except Exception as e:
            raise RuntimeError(f"文本转Word失败: {e}")
    
    @staticmethod
    def ofd_to_pdf(src: str, dst: str, log_func=None):
        """OFD 转 PDF - 使用 LibreOffice 或提示用户"""
        try:
            if log_func:
                log_func("  OFD 转换需要 LibreOffice 支持...")
            
            # OFD 转换策略：尝试使用 LibreOffice
            soffice_paths = [
                # macOS
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                # Windows
                "C:/Program Files/LibreOffice/program/soffice.exe",
                "C:/Program Files (x86)/LibreOffice/program/soffice.exe",
                # Linux/通用
                "soffice",
                "libreoffice"
            ]
            soffice = None
            for path in soffice_paths:
                if shutil.which(path) or os.path.exists(path):
                    soffice = path
                    break
            
            if soffice:
                import subprocess
                if log_func:
                    log_func("  使用 LibreOffice 转换...")
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf", "--outdir",
                     os.path.dirname(dst), src],
                    capture_output=True, timeout=120
                )
                lo_out = os.path.join(os.path.dirname(dst), Path(src).stem + ".pdf")
                if os.path.exists(lo_out) and os.path.abspath(lo_out) != os.path.abspath(dst):
                    shutil.move(lo_out, dst)
                if os.path.exists(dst):
                    return True
            
            # 如果 LibreOffice 不可用，尝试使用 PyMuPDF (部分支持)
            if log_func:
                log_func("  尝试使用 PyMuPDF 转换...")
            
            import fitz
            doc = fitz.open(src)
            doc.save(dst)
            doc.close()
            return True
            
        except Exception as e:
            raise RuntimeError(f"OFD转PDF失败: {e}。请安装 LibreOffice 以获得更好的 OFD 支持。")
    
    @staticmethod
    def _ofd_to_word(src: str, dst: str, log_func=None):
        """OFD 转 Word（中间步骤：先 OFD→PDF，再 PDF→Word）"""
        if log_func:
            log_func("  步骤1: OFD → PDF")
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            tmp_pdf = f.name
        try:
            ConversionEngine.ofd_to_pdf(src, tmp_pdf, log_func)
            if log_func:
                log_func("  步骤2: PDF → Word")
            ConversionEngine.pdf_to_word(tmp_pdf, dst, log_func)
            return True
        finally:
            if os.path.exists(tmp_pdf):
                os.unlink(tmp_pdf)
    
    @staticmethod
    def _ofd_to_txt(src: str, dst: str, log_func=None):
        """OFD 转文本（中间步骤：先 OFD→PDF，再 PDF→文本）"""
        if log_func:
            log_func("  步骤1: OFD → PDF")
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            tmp_pdf = f.name
        try:
            ConversionEngine.ofd_to_pdf(src, tmp_pdf, log_func)
            if log_func:
                log_func("  步骤2: PDF → 文本")
            ConversionEngine.pdf_to_txt(tmp_pdf, dst, log_func)
            return True
        finally:
            if os.path.exists(tmp_pdf):
                os.unlink(tmp_pdf)
    
    @classmethod
    def convert(cls, src: str, dst: str, src_fmt: str, dst_fmt: str, log_func=None):
        """统一转换入口"""
        key = (src_fmt, dst_fmt)
        routes = {
            ("pdf", "Word (.docx)"): cls.pdf_to_word,
            ("pdf", "Word-OCR (.docx)"): lambda s, d, l: cls.pdf_to_word(s, d, l, use_ocr=True),
            ("pdf", "Excel (.xlsx)"): cls.pdf_to_excel,
            ("pdf", "PPT (.pptx)"): cls.pdf_to_pptx,
            ("pdf", "图片 (.png)"): lambda s, d, l: cls.pdf_to_image(s, d, "png", l),
            ("pdf", "文本 (.txt)"): cls.pdf_to_txt,
            ("docx", "PDF (.pdf)"): cls.word_to_pdf,
            ("docx", "文本 (.txt)"): cls.word_to_txt,
            ("docx", "HTML (.html)"): cls.word_to_html,
            ("xlsx", "PDF (.pdf)"): cls.excel_to_pdf,
            ("xlsx", "CSV (.csv)"): cls.excel_to_csv,
            ("xlsx", "文本 (.txt)"): cls.excel_to_txt,
            ("pptx", "PDF (.pdf)"): cls.pptx_to_pdf,
            ("pptx", "图片 (.png)"): cls.pptx_to_image,
            ("ofd", "PDF (.pdf)"): cls.ofd_to_pdf,  # OFD转PDF
            ("ofd", "Word (.docx)"): cls._ofd_to_word,
            ("ofd", "文本 (.txt)"): cls._ofd_to_txt,
            ("png", "PDF (.pdf)"): cls.image_to_pdf,
            ("png", "JPG (.jpg)"): cls.image_to_image,
            ("png", "WebP (.webp)"): cls.image_to_image,
            ("jpg", "PDF (.pdf)"): cls.image_to_pdf,
            ("jpg", "PNG (.png)"): cls.image_to_image,
            ("jpg", "WebP (.webp)"): cls.image_to_image,
            ("jpeg", "PDF (.pdf)"): cls.image_to_pdf,
            ("jpeg", "PNG (.png)"): cls.image_to_image,
            ("jpeg", "WebP (.webp)"): cls.image_to_image,
            ("txt", "PDF (.pdf)"): cls.txt_to_pdf,
            ("txt", "Word (.docx)"): cls.txt_to_word,
        }
        fn = routes.get(key)
        if fn is None:
            raise ValueError(f"暂不支持 {src_fmt} → {dst_fmt} 的转换")
        return fn(src, dst, log_func)


# ===========================================================
# Main application
# ===========================================================

class FileFlowApp:
    def __init__(self, root):
        self.root = root
        self.root.title("FileFlow Pro")
        self.root.geometry("750x600")
        self.root.configure(bg=C["bg"])
        self.root.minsize(650, 500)
        
        # 窗口居中
        self._center_window()
        
        # 数据
        self.src_files = []
        self.target_fmt = tk.StringVar(value="请先选择源文件")
        self.output_dir = tk.StringVar(value=str(Path.home() / "Desktop"))
        self.pdf_password = tk.StringVar()
        self.is_converting = False
        
        # 构建界面
        self._build_ui()
        
    def _center_window(self):
        """窗口居中显示"""
        self.root.update_idletasks()
        width = 750
        height = 600
        x = (self.root.winfo_screenwidth() // 2) - (width // 2)
        y = (self.root.winfo_screenheight() // 2) - (height // 2)
        self.root.geometry(f'{width}x{height}+{x}+{y}')
        
    def _build_ui(self):
        """构建用户界面"""
        # 主容器
        main = tk.Frame(self.root, bg=C["bg"])
        main.pack(fill="both", expand=True, padx=20, pady=20)
        
        # 标题
        title = tk.Label(main, text="FileFlow Pro", font=("SF Pro Display", 28, "bold"),
                        bg=C["bg"], fg=C["text"])
        title.pack(anchor="w", pady=(0, 20))
        
        # 内容区域 - 左右分栏
        content = tk.Frame(main, bg=C["bg"])
        content.pack(fill="both", expand=True)
        content.grid_columnconfigure(0, weight=3)
        content.grid_columnconfigure(1, weight=2)
        
        # 左侧面板 - 文件选择
        left = self._create_card(content)
        left.grid(row=0, column=0, sticky="nsew", padx=(0, 10))
        
        self._build_left_panel(left)
        
        # 右侧面板 - 设置
        right = self._create_card(content)
        right.grid(row=0, column=1, sticky="nsew")
        
        self._build_right_panel(right)
        
    def _create_card(self, parent):
        """创建卡片容器"""
        card = tk.Frame(parent, bg=C["surface"], bd=1, relief="solid",
                       highlightbackground=C["border"])
        return card
        
    def _build_left_panel(self, parent):
        """构建左侧面板"""
        # 上传区域
        upload_frame = tk.Frame(parent, bg=C["surface"], padx=20, pady=20)
        upload_frame.pack(fill="x")
        
        # Upload icon (use plain text, no emoji)
        icon = tk.Label(upload_frame, text="[ + ]", font=("SF Pro", 28, "bold"),
                       bg=C["surface"], fg=C["blue"])
        icon.pack()
        
        # 提示文字
        hint = tk.Label(upload_frame, text="点击选择文件",
                       font=("SF Pro", 13), bg=C["surface"], fg=C["text2"])
        hint.pack(pady=5)
        
        # 支持格式
        formats = tk.Label(upload_frame, 
                          text="支持: PDF · Word · Excel · PPT · OFD · 图片 · 文本",
                          font=("SF Pro", 10), bg=C["surface"], fg=C["placeholder"])
        formats.pack()
        
        # 选择文件按钮
        btn_frame = tk.Frame(upload_frame, bg=C["surface"])
        btn_frame.pack(pady=15)
        
        select_btn = tk.Button(btn_frame, text="选择文件", command=self._select_files,
                              bg=C["blue"], fg="white", font=("SF Pro", 12, "bold"),
                              padx=25, pady=10, cursor="hand2", relief="flat",
                              activebackground="#003bb5", activeforeground="white")
        select_btn.pack(side="left", padx=5)
        
        clear_btn = tk.Button(btn_frame, text="清空", command=self._clear_files,
                             bg=C["surface"], fg=C["red"], font=("SF Pro", 12),
                             padx=20, pady=10, cursor="hand2", relief="solid",
                             highlightbackground=C["border"],
                             activebackground="#ffe5e5")
        clear_btn.pack(side="left", padx=5)
        
        # 文件列表
        list_frame = tk.Frame(parent, bg=C["surface"], padx=15, pady=10)
        list_frame.pack(fill="both", expand=True)
        
        # 列表标题
        list_title = tk.Frame(list_frame, bg=C["surface"])
        list_title.pack(fill="x", pady=(0, 5))
        
        tk.Label(list_title, text="已选择文件", font=("SF Pro", 12, "bold"),
                bg=C["surface"], fg=C["text"]).pack(side="left")
        
        self.file_count_label = tk.Label(list_title, text="0 个文件",
                                        font=("SF Pro", 11), bg=C["surface"],
                                        fg=C["placeholder"])
        self.file_count_label.pack(side="right")
        
        # 文件列表框
        self.file_listbox = tk.Listbox(list_frame, font=("SF Pro", 11),
                                      bg=C["surface"], fg=C["text"],
                                      selectbackground=C["blue_light"],
                                      highlightthickness=0, bd=1,
                                      relief="solid", height=10)
        self.file_listbox.pack(fill="both", expand=True)
        
        # 滚动条
        scrollbar = tk.Scrollbar(self.file_listbox)
        scrollbar.pack(side="right", fill="y")
        self.file_listbox.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.file_listbox.yview)
        
    def _build_right_panel(self, parent):
        """构建右侧面板"""
        parent.configure(padx=15, pady=15)
        
        # 转换设置标题
        tk.Label(parent, text="转换设置", font=("SF Pro", 16, "bold"),
                bg=C["surface"], fg=C["text"]).pack(anchor="w", pady=(0, 15))
        
        # 目标格式
        fmt_frame = tk.Frame(parent, bg=C["surface"])
        fmt_frame.pack(fill="x", pady=10)
        
        tk.Label(fmt_frame, text="转换为", font=("SF Pro", 12),
                bg=C["surface"], fg=C["text2"]).pack(anchor="w")
        
        # 使用自定义下拉框替代 ttk.Combobox
        self._create_custom_combobox(fmt_frame)
        
        # PDF密码
        pwd_frame = tk.Frame(parent, bg=C["surface"])
        pwd_frame.pack(fill="x", pady=10)
        
        tk.Label(pwd_frame, text="PDF密码（可选）", font=("SF Pro", 12),
                bg=C["surface"], fg=C["text2"]).pack(anchor="w")
        
        pwd_entry = tk.Entry(pwd_frame, textvariable=self.pdf_password,
                            font=("SF Pro", 11), show="*",
                            bg=C["surface"], relief="solid", bd=1)
        pwd_entry.pack(fill="x", pady=5)
        
        # 保存位置
        tk.Label(parent, text="保存位置", font=("SF Pro", 12),
                bg=C["surface"], fg=C["text2"]).pack(anchor="w", pady=(10, 0))
        
        dir_frame = tk.Frame(parent, bg=C["surface"])
        dir_frame.pack(fill="x", pady=5)
        
        dir_entry = tk.Entry(dir_frame, textvariable=self.output_dir,
                            font=("SF Pro", 11), bg=C["surface"],
                            relief="solid", bd=1)
        dir_entry.pack(side="left", fill="x", expand=True)
        
        browse_btn = tk.Button(dir_frame, text="浏览", command=self._browse_dir,
                              bg="#E5E5EA", font=("SF Pro", 10),
                              padx=12, pady=5, cursor="hand2",
                              relief="flat", activebackground="#D1D1D6")
        browse_btn.pack(side="right", padx=(5, 0))
        
        # 转换按钮
        self.convert_btn = tk.Button(parent, text="开始转换", command=self._start_convert,
                                    bg=C["blue"], fg="white",
                                    font=("SF Pro", 14, "bold"),
                                    padx=30, pady=14, cursor="hand2",
                                    relief="flat",
                                    activebackground="#003bb5")
        self.convert_btn.pack(fill="x", pady=20)
        
        # 日志区域
        tk.Label(parent, text="转换日志", font=("SF Pro", 12),
                bg=C["surface"], fg=C["text2"]).pack(anchor="w")
        
        self.log_text = tk.Text(parent, font=("SF Pro", 10), height=8,
                               bg=C["surface"], fg=C["text2"],
                               relief="solid", bd=1, wrap="word")
        self.log_text.pack(fill="both", expand=True, pady=5)
        
    def _create_custom_combobox(self, parent):
        """创建自定义下拉框，避免 ttk.Combobox 的问题"""
        self.combo_frame = tk.Frame(parent, bg=C["surface"])
        self.combo_frame.pack(fill="x", pady=5)
        
        # 显示当前选择的标签
        self.combo_display = tk.Label(self.combo_frame, 
                                     textvariable=self.target_fmt,
                                     font=("SF Pro", 11),
                                     bg="white", fg=C["text"],
                                     relief="solid", bd=1,
                                     padx=8, pady=6,
                                     anchor="w", cursor="hand2")
        self.combo_display.pack(fill="x")
        self.combo_display.bind("<Button-1>", lambda e: self._show_format_menu())
        
        # 下拉箭头
        self.arrow_label = tk.Label(self.combo_display, text="▼", 
                        font=("SF Pro", 10), bg="white", fg=C["placeholder"], cursor="hand2")
        self.arrow_label.place(relx=1.0, rely=0.5, anchor="e", x=-8)
        self.arrow_label.bind("<Button-1>", lambda e: self._show_format_menu())
        
        # 存储下拉选项
        self.format_options = []
        self._format_menu = None
        self._menu_open = False
        
    def _show_format_menu(self):
        """显示格式选择菜单 - 使用 tk.Menu 替代 Toplevel"""
        if not self.format_options:
            return
            
        if self._menu_open:
            return
            
        self._menu_open = True
        
        # 创建 tk.Menu 菜单
        self._format_menu = tk.Menu(self.root, tearoff=0, bg="white", 
                                    fg=C["text"], font=("SF Pro", 11),
                                    activebackground=C["blue_light"],
                                    activeforeground=C["text"],
                                    relief="solid", bd=1)
        
        # 添加菜单项
        for opt in self.format_options:
            self._format_menu.add_command(label=opt, 
                                         command=lambda o=opt: self._select_format(o))
        
        # 计算菜单位置
        x = self.combo_display.winfo_rootx()
        y = self.combo_display.winfo_rooty() + self.combo_display.winfo_height()
        
        # 显示菜单
        try:
            self._format_menu.post(x, y)
        except Exception as e:
            print(f"[ERROR] Menu display error: {e}")
            self._menu_open = False
            return
        
        # 绑定点击外部关闭菜单
        self.root.after(100, self._bind_menu_close)
        
    def _select_format(self, option):
        """选择格式"""
        self.target_fmt.set(option)
        self._close_format_menu()
        
    def _close_format_menu(self):
        """关闭格式菜单"""
        if self._format_menu:
            try:
                self._format_menu.unpost()
            except:
                pass
            try:
                self._format_menu.destroy()
            except:
                pass
            self._format_menu = None
        self._menu_open = False
        
    def _bind_menu_close(self):
        """绑定关闭菜单的事件"""
        if self._menu_open and self._format_menu:
            # 绑定全局点击事件
            self.root.bind_all("<Button-1>", self._on_global_click, add="+")
            
    def _on_global_click(self, event):
        """全局点击事件处理"""
        if not self._menu_open:
            return
            
        # 检查点击位置是否在组合框内
        widget = event.widget
        try:
            # 如果点击的是组合框相关组件，不关闭菜单
            if widget == self.combo_display or widget == self.arrow_label:
                return
            # 如果点击在组合框内
            if widget.winfo_toplevel() == self.combo_frame.winfo_toplevel():
                # 检查是否是组合框的子组件
                parent = widget
                while parent:
                    if parent == self.combo_frame:
                        return
                    try:
                        parent = parent.winfo_parent()
                        if parent:
                            parent = self.root.nametowidget(parent)
                    except:
                        break
        except:
            pass
            
        # 关闭菜单
        self._close_format_menu()
        # 解绑全局事件
        self.root.unbind_all("<Button-1>")
        
    def _select_files(self):
        """选择文件 - 使用 macOS 原生文件选择器"""
        try:
            self._log("正在打开文件选择器...")
            
            # 尝试使用 macOS 原生选择器
            result = select_files_native(
                title="选择要转换的文件",
                multiple=True,
                file_types=SUPPORTED_EXTS
            )
            
            # 如果原生选择器失败，使用手动输入
            if result is None:
                self._log("使用手动输入方式...")
                self._manual_file_input()
                return
                
            if result:
                self.src_files = result
                self._update_file_list()
                self._update_format_options()
                self._log(f"[OK] {len(self.src_files)} file(s) selected")
            else:
                self._log("No files selected")
                
        except Exception as e:
            self._log(f"文件选择错误: {e}")
            self._manual_file_input()
            
    def _manual_file_input(self):
        """手动输入文件路径"""
        dialog = tk.Toplevel(self.root)
        dialog.title("输入文件路径")
        dialog.geometry("550x350")
        dialog.transient(self.root)
        dialog.grab_set()
        dialog.configure(bg=C["bg"])
        
        # 居中显示
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (275)
        y = (dialog.winfo_screenheight() // 2) - (175)
        dialog.geometry(f"+{x}+{y}")
        
        tk.Label(dialog, text="请输入文件路径（每行一个）:", 
                font=("SF Pro", 12), bg=C["bg"], fg=C["text"]).pack(pady=15)
        
        text = tk.Text(dialog, font=("SF Pro", 11), height=12,
                      bg="white", fg=C["text"], relief="solid", bd=1)
        text.pack(fill="both", expand=True, padx=20, pady=5)
        
        btn_frame = tk.Frame(dialog, bg=C["bg"])
        btn_frame.pack(pady=15)
        
        def confirm():
            paths = [p.strip() for p in text.get("1.0", tk.END).strip().split("\n") if p.strip()]
            valid_paths = [p for p in paths if os.path.isfile(p)]
            if valid_paths:
                self.src_files = valid_paths
                self._update_file_list()
                self._update_format_options()
                self._log(f"[OK] {len(self.src_files)} file(s) selected")
            else:
                self._log("No valid files found")
            dialog.destroy()
            
        def cancel():
            dialog.destroy()
            
        tk.Button(btn_frame, text="确定", command=confirm,
                 bg=C["blue"], fg="white", font=("SF Pro", 11),
                 padx=25, pady=8, relief="flat",
                 activebackground="#003bb5").pack(side="left", padx=5)
                 
        tk.Button(btn_frame, text="取消", command=cancel,
                 bg="#E5E5EA", fg=C["text"], font=("SF Pro", 11),
                 padx=25, pady=8, relief="flat",
                 activebackground="#D1D1D6").pack(side="left", padx=5)
        
    def _update_file_list(self):
        """更新文件列表显示"""
        self.file_listbox.delete(0, tk.END)
        for f in self.src_files:
            display_name = f"  {Path(f).name}"
            self.file_listbox.insert(tk.END, display_name)
        self.file_count_label.config(text=f"{len(self.src_files)} 个文件")
        
    def _update_format_options(self):
        """根据源文件类型更新格式选项"""
        if not self.src_files:
            self.format_options = []
            self.target_fmt.set("请先选择源文件")
            return
            
        # 获取第一个文件的扩展名
        ext = Path(self.src_files[0]).suffix.lower().lstrip('.')
        
        self.format_options = CONVERT_OPTIONS.get(ext, [])
        
        if self.format_options:
            self.target_fmt.set(self.format_options[0])
            self._log(f"可用格式: {', '.join(self.format_options)}")
        else:
            self.target_fmt.set("暂不支持此格式")
            self.format_options = []
            self._log(f"警告: 不支持 {ext} 格式")
            
    def _clear_files(self):
        """清空文件列表"""
        self.src_files = []
        self._update_file_list()
        self.target_fmt.set("请先选择源文件")
        self.format_options = []
        self._log("已清空文件列表")
        
    def _browse_dir(self):
        """浏览输出目录 - 使用 macOS 原生文件夹选择器"""
        try:
            self._log("正在打开文件夹选择器...")
            
            result = select_folder_native(
                title="选择保存位置",
                initial_dir=self.output_dir.get()
            )
            
            if result is None:
                # 使用手动输入
                self._manual_dir_input()
                return
                
            if result:
                self.output_dir.set(result)
                self._log(f"[OK] Save location: {result}")
            else:
                self._log("No folder selected")
                
        except Exception as e:
            self._log(f"文件夹选择错误: {e}")
            self._manual_dir_input()
            
    def _manual_dir_input(self):
        """手动输入目录路径"""
        dialog = tk.Toplevel(self.root)
        dialog.title("输入保存路径")
        dialog.geometry("500x150")
        dialog.transient(self.root)
        dialog.grab_set()
        dialog.configure(bg=C["bg"])
        
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (250)
        y = (dialog.winfo_screenheight() // 2) - (75)
        dialog.geometry(f"+{x}+{y}")
        
        tk.Label(dialog, text="请输入保存文件夹路径:", 
                font=("SF Pro", 12), bg=C["bg"], fg=C["text"]).pack(pady=10)
        
        entry = tk.Entry(dialog, font=("SF Pro", 11), width=50,
                        bg="white", relief="solid", bd=1)
        entry.pack(padx=20, pady=5)
        entry.insert(0, self.output_dir.get())
        
        btn_frame = tk.Frame(dialog, bg=C["bg"])
        btn_frame.pack(pady=10)
        
        def confirm():
            path = entry.get().strip()
            if path and os.path.isdir(path):
                self.output_dir.set(path)
                self._log(f"[OK] Save location: {path}")
            else:
                self._log("错误: 无效的文件夹路径")
            dialog.destroy()
            
        def cancel():
            dialog.destroy()
            
        tk.Button(btn_frame, text="确定", command=confirm,
                 bg=C["blue"], fg="white", font=("SF Pro", 11),
                 padx=20, pady=6, relief="flat").pack(side="left", padx=5)
                 
        tk.Button(btn_frame, text="取消", command=cancel,
                 bg="#E5E5EA", fg=C["text"], font=("SF Pro", 11),
                 padx=20, pady=6, relief="flat").pack(side="left", padx=5)
            
    def _log(self, msg):
        """添加日志（线程安全）"""
        timestamp = time.strftime("%H:%M:%S")
        def _insert():
            self.log_text.insert(tk.END, f"[{timestamp}] {msg}\n")
            self.log_text.see(tk.END)
        try:
            self.root.after(0, _insert)
        except RuntimeError:
            # 窗口已关闭时忽略
            pass
        
    def _start_convert(self):
        """开始转换"""
        if not self.src_files:
            messagebox.showwarning("提示", "请先选择要转换的文件")
            return
            
        if self.is_converting:
            messagebox.showwarning("提示", "转换正在进行中，请稍候")
            return
            
        target = self.target_fmt.get()
        if "请先选择" in target or "暂不支持" in target:
            messagebox.showwarning("提示", "请选择有效的目标格式")
            return
            
        # 检查输出目录
        out_dir = self.output_dir.get()
        if not os.path.isdir(out_dir):
            messagebox.showerror("错误", f"输出目录不存在: {out_dir}")
            return
            
        # 在后台线程执行转换
        self.is_converting = True
        self.convert_btn.config(state="disabled", text="转换中...")
        self._log(f"=" * 40)
        self._log(f"开始转换 {len(self.src_files)} 个文件")
        self._log(f"目标格式: {target}")
        self._log(f"输出目录: {out_dir}")
        self._log(f"=" * 40)
        
        thread = threading.Thread(target=self._convert_worker, args=(target, out_dir))
        thread.daemon = True
        thread.start()
        
    def _convert_worker(self, target, out_dir):
        """转换工作线程"""
        try:
            success_count = 0
            failed_files = []
            
            for i, src in enumerate(self.src_files, 1):
                self._log(f"[{i}/{len(self.src_files)}] 处理: {Path(src).name}")
                
                try:
                    # 获取源文件扩展名
                    src_ext = Path(src).suffix.lower().lstrip('.')
                    
                    # 获取输出扩展名
                    out_ext = EXT_MAP.get(target, ".out")
                    out_name = f"{Path(src).stem}_converted{out_ext}"
                    out_path = os.path.join(out_dir, out_name)
                    
                    # 检查源文件是否存在
                    if not os.path.exists(src):
                        raise FileNotFoundError(f"源文件不存在: {src}")
                    
                    # 检查输出目录是否可写
                    if not os.access(out_dir, os.W_OK):
                        raise PermissionError(f"输出目录没有写入权限: {out_dir}")
                    
                    # 执行转换
                    ConversionEngine.convert(src, out_path, src_ext, target, self._log)
                    
                    # 验证输出文件是否生成
                    if os.path.exists(out_path):
                        file_size = os.path.getsize(out_path)
                        self._log(f"  [OK] Done: {out_name} ({file_size} bytes)")
                        success_count += 1
                    else:
                        raise RuntimeError("转换后未找到输出文件")
                    
                except Exception as e:
                    error_msg = str(e)
                    import traceback
                    tb_str = traceback.format_exc()
                    self._log(f"  [FAIL] Failed: {error_msg}")
                    # 记录详细错误信息到日志（仅在开发时显示）
                    for line in tb_str.split('\n')[:5]:  # 只显示前5行
                        if line.strip():
                            self._log(f"    {line}")
                    failed_files.append((Path(src).name, error_msg))
            
            self.root.after(0, self._convert_done, success_count, failed_files)
            
        except Exception as e:
            import traceback
            tb_str = traceback.format_exc()
            self.root.after(0, self._convert_done, 0, [("未知错误", f"{str(e)}\n{tb_str}")])
            
    def _convert_done(self, success_count, failed_files):
        """转换完成回调"""
        self.is_converting = False
        self.convert_btn.config(state="normal", text="开始转换")
        
        self._log(f"=" * 40)
        total = len(self.src_files)
        if failed_files:
            self._log(f"完成: {success_count}/{total} 成功, {len(failed_files)} 失败")
            for name, err in failed_files:
                self._log(f"  - {name}: {err}")
            messagebox.showwarning("部分失败", 
                f"成功: {success_count} 个\n失败: {len(failed_files)} 个\n\n请查看日志了解详情")
        else:
            self._log(f"[OK] All {success_count} file(s) converted successfully!")
            messagebox.showinfo("完成", f"成功转换 {success_count} 个文件！\n\n文件保存在:\n{self.output_dir.get()}")
        self._log(f"=" * 40)


# ===========================================================
# Entry point
# ===========================================================

if __name__ == "__main__":
    try:
        root = tk.Tk()
        app = FileFlowApp(root)
        root.mainloop()
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        temp_dir = os.environ.get('TEMP', os.environ.get('TMP', os.getcwd()))
        # Always print to console (visible when console=True)
        print(f"[FATAL] {tb}", file=sys.stderr)
        sys.stderr.flush()
        try:
            import tkinter.messagebox as mb
            mb.showerror("FileFlow Pro - Fatal Error", f"{e}\n\n{tb}")
        except Exception:
            log_path = os.path.join(temp_dir, 'fileflowpro_error.log')
            try:
                with open(log_path, 'w', encoding='utf-8') as f:
                    f.write(f"FileFlow Pro Fatal Error:\n{tb}\n")
                print(f"[FATAL] Log written to: {log_path}", file=sys.stderr)
            except Exception:
                pass
        _pause_on_error(f"FATAL ERROR.\nLog: {os.path.join(temp_dir, 'fileflowpro_error.log')}\nPress Enter to exit...")
        sys.exit(1)
